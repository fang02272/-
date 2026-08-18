# -*- coding: utf-8 -*-
"""生成本周汇报 PPT v2 — 真表格/柱状图/架构图，商务风格"""
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ---- 商务配色 ----
NAVY = RGBColor(0x1F, 0x3A, 0x5F)    # 深蓝(主)
BLUE = RGBColor(0x2E, 0x74, 0xB5)    # 中蓝
LBLUE = RGBColor(0xD6, 0xE4, 0xF0)   # 浅蓝
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x2B, 0x2B, 0x2B)
GRAY = RGBColor(0x6B, 0x6B, 0x6B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xC5, 0x5A, 0x11)
BG = RGBColor(0xF4, 0xF7, 0xFA)
FONT = "微软雅黑"

SW, SH = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = SW, SH


def _setfont(r, size, bold=False, color=DARK):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def title_bar(s, text, num=None):
    # 顶部标题条
    bar = s.shapes.add_shape(1, 0, 0, SW, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tf = bar.text_frame; tf.margin_left = Inches(0.5); tf.margin_top = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text; _setfont(r, 26, True, WHITE)
    # 页脚页码
    foot = s.shapes.add_textbox(SW - Inches(1.0), SH - Inches(0.5), Inches(0.8), Inches(0.4))
    ftf = foot.text_frame; fp = ftf.paragraphs[0]; fp.alignment = PP_ALIGN.RIGHT
    fr = fp.add_run(); fr.text = str(num or ""); _setfont(fr, 12, False, GRAY)


def add_text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT):
    """runs: list of paragraphs; each paragraph = list of (text, size, bold, color)"""
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for para in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align; p.space_after = Pt(5)
        for text, sz, bd, cl in para:
            r = p.add_run(); r.text = text; _setfont(r, sz, bd, cl)
    return tb


def add_box(s, x, y, w, h, fill, line=None):
    b = s.shapes.add_shape(1, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill
    if line:
        b.line.color.rgb = line; b.line.width = Pt(1)
    else:
        b.line.fill.background()
    return b


def add_table(s, x, y, w, h, data, col_widths=None, header_fill=NAVY, font_size=13):
    """data: 二维列表，首行表头"""
    rows, cols = len(data), len(data[0])
    shape = s.shapes.add_table(rows, cols, x, y, w, h)
    tbl = shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for ri in range(rows):
        for ci in range(cols):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(data[ri][ci])
            if ri == 0:
                _setfont(r, font_size, True, WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                bold = (ci == 0)
                _setfont(r, font_size, bold, DARK if not bold else NAVY)
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else LBLUE
    return tbl


def add_bar_chart(s, x, y, w, h, title, categories, series, colors=None):
    cd = CategoryChartData()
    cd.categories = categories
    for name, vals in series:
        cd.add_series(name, vals)
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, cd)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    try:
        chart.legend.font.size = Pt(11); chart.legend.font.name = FONT
    except Exception:
        pass
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    try:
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(13); run.font.bold = True; run.font.name = FONT
    except Exception:
        pass
    plot = chart.plots[0]
    plot.gap_width = 80
    for i, s_ in enumerate(plot.series):
        s_.format.fill.solid()
        c = colors[i] if colors else BLUE
        s_.format.fill.fore_color.rgb = c
    return chart


# ================= P1 封面 =================
s = add_slide()
add_box(s, 0, 0, SW, SH, NAVY)
add_box(s, 0, Inches(5.2), SW, Inches(0.06), ORANGE)
add_text(s, Inches(1), Inches(1.4), Inches(11.3), Inches(0.5),
         [[("国家科技重大专项", 18, False, LBLUE)]], align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.1), Inches(11.3), Inches(1.4),
         [[("焊接工艺大模型知识问答系统", 42, True, WHITE)]], align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(3.6), Inches(11.3), Inches(0.6),
         [[("—— 机器人焊接能力升级（本周汇报）——", 20, False, LBLUE)]], align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.6), Inches(11.3), Inches(0.6),
         [[("汇报人：李芳      时间：2026.08.15", 16, False, WHITE)]], align=PP_ALIGN.CENTER)

# ================= P2 目录 =================
s = add_slide(); title_bar(s, "目录", 2)
items = [("01", "本周工作概览", "机器人焊接能力落地"),
         ("02", "具体进展（基线对比 · 数据 · 成果）", "有基线 / 有数据 / 有成果"),
         ("03", "需求 · 突破 · 形态 · 指标", "4个维度"),
         ("04", "下一步计划", "知识库建立与优化")]
y = Inches(1.7)
for num, t, sub in items:
    add_box(s, Inches(1.2), y, Inches(11), Inches(1.1), WHITE, LBLUE)
    add_text(s, Inches(1.6), y + Inches(0.2), Inches(1.2), Inches(0.7), [[(num, 28, True, ORANGE)]])
    add_text(s, Inches(3.0), y + Inches(0.18), Inches(7), Inches(0.6), [[(t, 20, True, NAVY)]])
    add_text(s, Inches(3.0), y + Inches(0.62), Inches(7), Inches(0.4), [[(sub, 13, False, GRAY)]])
    y += Inches(1.35)

# ================= P3 技术路线（架构图）=================
s = add_slide(); title_bar(s, "技术路线：四步闭环", 3)
add_text(s, Inches(1), Inches(1.2), Inches(11), Inches(0.5),
         [[("查参数 → 出机器可读工艺卡片 → 指导机器人仿真/路径规划", 17, True, NAVY)]])
boxes = [
    ("① 知识获取", "书 / 手册\n期刊 / 标准\n技术合同", Inches(0.8)),
    ("② 知识解析", "六阶段关键词\njieba 分词\n表格重建", Inches(3.8)),
    ("③ 知识应用", "意图路由\n工艺卡片\n机器人字段", Inches(6.8)),
    ("④ 持续优化", "案例 / 实测\n卡诺普参数\n校准", Inches(9.8)),
]
for t, sub, x in boxes:
    add_box(s, x, Inches(2.2), Inches(2.6), Inches(2.2), LBLUE, BLUE)
    add_text(s, x, Inches(2.4), Inches(2.6), Inches(0.6), [[(t, 17, True, NAVY)]], align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.1), Inches(2.6), Inches(1.0), [[(ln, 13, False, DARK)] for ln in sub.split("\n")],
             align=PP_ALIGN.CENTER)
# 箭头
for ax in (Inches(3.5), Inches(6.5), Inches(9.5)):
    ar = s.shapes.add_shape(24, ax, Inches(3.0), Inches(0.4), Inches(0.5))  # right arrow
    ar.fill.solid(); ar.fill.fore_color.rgb = ORANGE; ar.line.fill.background()
# 闭环回箭头
ar = s.shapes.add_shape(24, Inches(1.5), Inches(4.8), Inches(10.3), Inches(0.45))  # right arrow 底部
ar.fill.solid(); ar.fill.fore_color.rgb = GREEN; ar.line.fill.background()
add_text(s, Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
         [[("④ 持续优化：案例 + 实测 + 卡诺普参数校准 → 反哺知识库（闭环）", 14, False, GREEN)]],
         align=PP_ALIGN.CENTER)

# ================= P4 有基线：上周→本周（真表格）=================
s = add_slide(); title_bar(s, "有基线：上周 → 本周对比", 4)
data = [
    ["维度", "上周（基线）", "本周", "提升"],
    ["词库", "1317 关键词 / 123 同义词组", "1317 / 123（保持）+ 机器人术语", "稳定"],
    ["专家库", "123 概念", "123 概念 + 机器人字段", "能力扩展"],
    ["检索验证", "159 行向量，6 类查询通过", "同上 + 管道/船型焊场景", "场景扩展"],
    ["工艺卡片", "手工焊参数（焊条电弧焊）", "机器人焊接卡片（GMAW/焊丝/层道/姿态/管道）", "核心突破"],
]
add_table(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.6), data,
          col_widths=[Inches(1.8), Inches(3.6), Inches(4.2), Inches(2.5)])
# 柱状图：关键词提升
add_text(s, Inches(0.6), Inches(4.4), Inches(6), Inches(0.5),
         [[("关键词扩充对比（词库建设）", 15, True, NAVY)]])
add_bar_chart(s, Inches(0.6), Inches(4.9), Inches(6), Inches(2.2), "",
              ["材料焊接原理", "焊接结构原理"],
              [("上周", [278, 114]), ("本周", [1116, 123])],
              colors=[GRAY, GREEN])
# 表格数对比
add_text(s, Inches(7.2), Inches(4.4), Inches(5), Inches(0.5),
         [[("表格重建（扫描版结构化）", 15, True, NAVY)]])
add_bar_chart(s, Inches(7.2), Inches(4.9), Inches(5), Inches(2.2), "",
              ["实用焊接手册"],
              [("上周", [0]), ("本周", [112])],
              colors=[GRAY, BLUE])

# ================= P5 有数据：机器人焊接能力 =================
s = add_slide(); title_bar(s, "有数据：机器人焊接能力", 5)
# 左：关键数据卡
add_text(s, Inches(0.6), Inches(1.2), Inches(7), Inches(0.5),
         [[("按母材自动推荐（示例 Q345 12mm 对接）", 16, True, NAVY)]])
data_cards = [
    ("默认工艺", "SMAW → GMAW/MIG", "人工焊 → 机器人焊丝工艺"),
    ("焊丝选型", "ER50-6 · 1.2mm", "低碳钢→ER50-6、不锈钢→ER308L"),
    ("保护气", "80%Ar + 20%CO₂", "MAG 富氩混合气"),
    ("TCP/导电嘴", "导电嘴到工件 12-15mm", "干伸长 15-18mm"),
    ("枪姿态", "工作角75-85° / 行走角10-15°", "绕枪轴0°"),
    ("船型焊", "PA/1G 位置", "坡口转45°朝上，熔池水平最稳"),
]
yy = Inches(1.8)
for t, v, sub in data_cards:
    add_box(s, Inches(0.6), yy, Inches(6.2), Inches(0.75), WHITE, LBLUE)
    add_text(s, Inches(0.8), yy + Inches(0.06), Inches(1.6), Inches(0.6), [[(t, 14, True, NAVY)]])
    add_text(s, Inches(2.4), yy + Inches(0.06), Inches(4.2), Inches(0.6),
             [[(v, 14, True, ORANGE), ("   " + sub, 11, False, GRAY)]])
    yy += Inches(0.85)
# 右：层道序列真表格
add_text(s, Inches(7.2), Inches(1.2), Inches(5.5), Inches(0.5),
         [[("层道序列：打底 → 填充 → 盖面（具体电流电压）", 16, True, NAVY)]])
layer_data = [
    ["焊道", "电流", "电压", "焊速", "送丝", "焊宽"],
    ["打底", "120-140A", "18-20V", "25-35", "6-8", "4-5mm"],
    ["填充", "160-190A", "21-24V", "28-35", "7-9", "6-7mm"],
    ["盖面", "180-210A", "23-25V", "25-30", "8-10", "7-9mm"],
]
add_table(s, Inches(7.2), Inches(1.8), Inches(5.5), Inches(2.4), layer_data,
          col_widths=[Inches(1.0), Inches(1.3), Inches(1.1), Inches(0.9), Inches(0.7), Inches(0.9)], font_size=12)
add_text(s, Inches(7.2), Inches(4.4), Inches(5.5), Inches(2.2), [
    [("管道 6G 场景：", 14, True, NAVY)],
    [("· 固定管分两个半圈（6点→12点）", 13, False, DARK)],
    [("· 每半圈含斜仰/斜立/斜平", 13, False, DARK)],
    [("· 电流比平焊降 10-15%", 13, False, DARK)],
])

# ================= P6 有成果：工艺卡片升级 =================
s = add_slide(); title_bar(s, "有成果：工艺卡片升级（机器人可执行）", 6)
# 左：平板卡片示意
add_box(s, Inches(0.7), Inches(1.3), Inches(5.9), Inches(5.4), WHITE, BLUE)
add_box(s, Inches(0.7), Inches(1.3), Inches(5.9), Inches(0.6), BLUE)
add_text(s, Inches(0.9), Inches(1.4), Inches(5.5), Inches(0.5),
         [[("🦾 机器人焊接工艺卡片（Q345 12mm）", 15, True, WHITE)]])
add_text(s, Inches(1.0), Inches(2.1), Inches(5.4), Inches(4.3), [
    [("工艺：GMAW/MIG（熔化极氩弧焊）", 13, True, DARK)],
    [("焊丝：ER50-6 · 1.2mm ｜ 保护气：80%Ar+20%CO₂", 12, False, DARK)],
    [("", 6, False, DARK)],
    [("层道：", 13, True, NAVY)],
    [("  打底  120-140A / 18-20V", 12, False, DARK)],
    [("  填充  160-190A / 21-24V", 12, False, DARK)],
    [("  盖面  180-210A / 23-25V", 12, False, DARK)],
    [("", 6, False, DARK)],
    [("姿态：工作角75-85° ｜ 船型焊 PA", 12, False, DARK)],
    [("装备：MR2010_1 / APW50N / NBC-500RP", 11, False, GRAY)],
    [("", 6, False, DARK)],
    [("✔ 机器可读 JSON（仿真/机器人导入）", 12, True, GREEN)],
    [("✔ 打印 PDF / 复制 JSON", 12, True, GREEN)],
])
# 右：管道卡片示意
add_box(s, Inches(6.9), Inches(1.3), Inches(5.7), Inches(5.4), WHITE, ORANGE)
add_box(s, Inches(6.9), Inches(1.3), Inches(5.7), Inches(0.6), ORANGE)
add_text(s, Inches(7.1), Inches(1.4), Inches(5.3), Inches(0.5),
         [[("🛢 管道焊接工艺卡片（45°固定管 6G）", 15, True, WHITE)]])
add_text(s, Inches(7.2), Inches(2.1), Inches(5.2), Inches(4.3), [
    [("工艺：GMAW/MIG ｜ 12mm", 13, True, DARK)],
    [("焊丝：ER50-6 · 1.2mm", 12, False, DARK)],
    [("", 6, False, DARK)],
    [("管道策略：", 13, True, ORANGE)],
    [("  固定管 6点→12点半圈分段", 12, False, DARK)],
    [("  每半圈含斜仰/斜立/斜平", 12, False, DARK)],
    [("  电流比平焊降10-15%", 12, False, DARK)],
    [("", 6, False, DARK)],
    [("船型焊：PA位置", 13, True, NAVY)],
    [("  坡口转45°朝上，熔池水平", 12, False, DARK)],
    [("", 6, False, DARK)],
    [("✔ 前端「机器人焊接方案」板块展示", 12, True, GREEN)],
])
# 底部
add_box(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.55), LBLUE)
add_text(s, Inches(1.0), Inches(6.82), Inches(11.4), Inches(0.4),
         [[("前端展示：工艺卡片新增「机器人焊接方案」板块（层道电流电压表/管道策略）｜ 命令行 run_welding_qa.py 同步输出",
            12, False, NAVY)]])

# ================= P7 需求/突破/形态/指标 =================
s = add_slide(); title_bar(s, "有需求 · 有突破 · 有形态 · 有指标", 7)
quads = [
    ("有需求", "手册偏人工焊（焊条电弧焊/运条），参数是范围值，机器人用不上；需要「确定值 + 分层 + 姿态」的机器人焊接方案", NAVY),
    ("有突破", "默认工艺切 GMAW/MIG、新增机器人焊接规则库（robot_welding.py）、层道级电流电压、管道 6G/船型焊识别", BLUE),
    ("有形态", "工艺卡片 = 机器可读 JSON（供仿真/机器人）+ 前端可视化 + 打印 PDF", GREEN),
    ("有指标", "机器人场景识别 ✓ ｜ 层道序列生成 ✓ ｜ 管道6G分段 ✓ ｜ 测试套件全部通过 ✓", ORANGE),
]
pos = [(Inches(0.8), Inches(1.5)), (Inches(6.9), Inches(1.5)),
       (Inches(0.8), Inches(4.4)), (Inches(6.9), Inches(4.4))]
for (t, it, cl), (x, y) in zip(quads, pos):
    add_box(s, x, y, Inches(5.6), Inches(2.7), WHITE, cl)
    add_box(s, x, y, Inches(5.6), Inches(0.6), cl)
    add_text(s, x + Inches(0.2), y + Inches(0.05), Inches(5.2), Inches(0.5),
             [[(f"【{t}】", 17, True, WHITE)]])
    add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(5.2), Inches(1.8),
             [[(it, 14, False, DARK)]])

# ================= P8 下一步计划（简版）=================
s = add_slide(); title_bar(s, "下一步计划：知识库建立与优化", 8)
add_text(s, Inches(1), Inches(1.3), Inches(11.3), Inches(0.5),
         [[("持续知识喂入 → 知识库建立与优化（支撑机器人工艺卡片更贴合现场）", 17, True, NAVY)]])
# 三步：喂入 → 建立 → 优化
steps = [
    ("知识喂入", "期刊 / 标准\n技术合同\n卡诺普工艺参数", NAVY),
    ("知识库建立", "扫描件 GPU 识别\n表格重建入库\n多书统一管理", BLUE),
    ("知识库优化", "检索精度提升\n词库/同义词扩充\n卡片参数校准", GREEN),
]
x = Inches(0.8)
for t, sub, cl in steps:
    add_box(s, x, Inches(2.2), Inches(3.6), Inches(2.8), LBLUE, cl)
    add_text(s, x + Inches(0.2), Inches(2.4), Inches(3.2), Inches(0.6), [[(t, 18, True, cl)]])
    add_text(s, x + Inches(0.2), Inches(3.2), Inches(3.2), Inches(1.6),
             [[(ln, 13, False, DARK)] for ln in sub.split("\n")])
    if x < Inches(9.0):
        ar = s.shapes.add_shape(24, x + Inches(3.6), Inches(3.4), Inches(0.4), Inches(0.5))
        ar.fill.solid(); ar.fill.fore_color.rgb = ORANGE; ar.line.fill.background()
    x += Inches(4.0)
add_text(s, Inches(1), Inches(5.5), Inches(11.3), Inches(1.2), [
    [("重点：加入机器人实际参数（卡诺普），校准层道电流电压为确定值", 14, True, ORANGE)],
    [("从技术合同案例提炼焊接经验规则；期刊/标准导入补充场景知识", 13, False, GRAY)],
])

# ================= P9 感谢 =================
s = add_slide()
add_box(s, 0, 0, SW, SH, NAVY)
add_text(s, Inches(1), Inches(3.0), Inches(11.3), Inches(1.0),
         [[("请各位专家批评指正！", 38, True, WHITE)]], align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.3), Inches(11.3), Inches(0.6),
         [[("焊接工艺大模型知识问答系统", 16, False, LBLUE)]], align=PP_ALIGN.CENTER)

out = r"E:\PyCharmMisProject\docs\机器人焊接知识智能问答系统（08.15）v2.pptx"
prs.save(out)
print(f"✅ 已生成: {out}")
